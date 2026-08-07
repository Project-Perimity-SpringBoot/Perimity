"""
Builds the Perimity final-evaluation pitch deck.

Five slides, sized for a ten-minute slot that is mostly demo:
  2 min  problem + solution   -> slides 2 and 3 carry this
  5 min  live demonstration   -> the slides are off-screen for this
  3 min  Q&A                  -> slides 4 and 5 exist to answer questions

The slides are built to be READ IN THREE SECONDS, not narrated.

python-pptx cannot write animations, so the deck is assembled first and the
OOXML timing tree is injected afterwards - see animate() at the bottom.
"""

import re
import shutil
import zipfile
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ----------------------------------------------------------------- palette
# Taken from the product itself - the gate pass, the pass email and the app
# card are all this violet. Amber is the single sharp accent, used only where
# something must be noticed: the differentiator and the live status.
INK        = RGBColor(0x2E, 0x10, 0x65)
INK_LIFT   = RGBColor(0x3B, 0x16, 0x7A)   # for layered shapes on the dark bg
VIOLET     = RGBColor(0x4C, 0x1D, 0x95)
VIOLET_MID = RGBColor(0x6D, 0x28, 0xD9)
LILAC      = RGBColor(0xA7, 0x8B, 0xFA)
SOFT       = RGBColor(0xED, 0xE9, 0xFE)
PAPER      = RGBColor(0xF7, 0xF5, 0xFF)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
TEXT       = RGBColor(0x1E, 0x1B, 0x4B)
MUTED      = RGBColor(0x6B, 0x72, 0x80)
MUTED_DARK = RGBColor(0xC4, 0xB5, 0xFD)
AMBER      = RGBColor(0xF5, 0x9E, 0x0B)
AMBER_SOFT = RGBColor(0xFD, 0xE6, 0x8A)

HEAD = "Cambria"
BODY = "Calibri"

W, H = 13.333, 7.5

OUT = "Perimity-Pitch-v2.pptx"

# Shape ids to animate, per slide index. Filled while building.
timeline = {}


# --------------------------------------------------------------- utilities


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def flat(shp):
    """Kill the inherited theme shadow. Used for everything decorative."""
    shp.shadow.inherit = False
    return shp


def fill_bg(slide, colour):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(W), Inches(H))
    shp.fill.solid()
    shp.fill.fore_color.rgb = colour
    shp.line.fill.background()
    return flat(shp)


def card(slide, x, y, w, h, colour, radius=0.06, line=None, shadow=False):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shp.adjustments[0] = radius
    shp.fill.solid()
    shp.fill.fore_color.rgb = colour
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(1)
    if not shadow:
        flat(shp)
    return shp


def oval(slide, x, y, w, h, colour):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = colour
    shp.line.fill.background()
    return flat(shp)


def rect(slide, x, y, w, h, colour, shape=MSO_SHAPE.RECTANGLE):
    shp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = colour
    shp.line.fill.background()
    return flat(shp)


def text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space=0, line_spacing=None):
    """runs = [(string, size, bold, colour, font)] - one paragraph each."""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor

    for i, (s, size, bold, colour, font) in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if space and i:
            p.space_before = Pt(space)
        if line_spacing:
            p.line_spacing = line_spacing
        r = p.add_run()
        r.text = s
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = colour
        r.font.name = font
    return box


# ---------------------------------------------------------------- glyphs
# Small icons built from primitives. Real icon fonts are not available in this
# environment, and clip-art would cheapen the deck - two or three shapes read
# as an icon at this size and stay on-palette.


def glyph_person(slide, cx, cy, colour):
    return [oval(slide, cx - 0.075, cy - 0.135, 0.15, 0.15, colour),
            rect(slide, cx - 0.13, cy + 0.035, 0.26, 0.145, colour,
                 MSO_SHAPE.ROUNDED_RECTANGLE)]


def glyph_door(slide, cx, cy, colour):
    return [rect(slide, cx - 0.10, cy - 0.15, 0.20, 0.30, colour,
                 MSO_SHAPE.ROUNDED_RECTANGLE),
            oval(slide, cx + 0.035, cy - 0.02, 0.05, 0.05, WHITE)]


def glyph_scan(slide, cx, cy, colour):
    return [rect(slide, cx + dx, cy + dy, 0.09, 0.09, colour)
            for dx, dy in ((-0.12, -0.13), (0.03, -0.13), (-0.12, 0.02), (0.03, 0.02))]


def glyph_log(slide, cx, cy, colour):
    return [rect(slide, cx - 0.13, cy - 0.13 + i * 0.105, 0.26, 0.05, colour,
                 MSO_SHAPE.ROUNDED_RECTANGLE) for i in range(3)]


def ids(*shapes):
    """Flatten shapes and lists of shapes into a list of shape ids."""
    out = []
    for item in shapes:
        if isinstance(item, (list, tuple)):
            out.extend(sh.shape_id for sh in item)
        else:
            out.append(item.shape_id)
    return out


# ------------------------------------------------------------------ deck

prs = Presentation()
prs.slide_width = Inches(W)
prs.slide_height = Inches(H)


# =========================================================== 1 - title
s = blank(prs)
fill_bg(s, INK)

# Layered discs behind the code block. They give the flat background depth
# without a gradient, which pptx cannot write natively.
oval(s, 7.3, 0.55, 6.2, 6.2, INK_LIFT)
oval(s, 8.05, 1.35, 4.7, 4.7, INK)

QX, QY, CELL = 8.62, 1.95, 0.265
PATTERN = [
    "11111011011111",
    "10001001010001",
    "10111010110111",
    "10111000010111",
    "10111011010111",
    "10001010010001",
    "11111010111111",
    "00000110100000",
    "10110101101101",
    "01001011010010",
    "11011000101101",
    "00000101011011",
    "11111010110101",
    "10001011001110",
]
for r, row in enumerate(PATTERN):
    for c, ch in enumerate(row):
        if ch != "1":
            continue
        finder = (r < 7 and c < 7) or (r < 7 and c > 6) or (r > 6 and c < 7)
        shade = WHITE if finder else LILAC
        rect(s, QX + c * CELL, QY + r * CELL, CELL * 0.84, CELL * 0.84, shade)

t1 = text(s, 0.95, 2.15, 7.2, 1.2, [("PERIMITY", 60, True, WHITE, HEAD)])
t2 = text(s, 0.95, 3.32, 7.2, 1.5,
          [("Smart Campus Access &\nGate Pass Management", 23, False, SOFT, BODY)],
          line_spacing=1.25)
t3 = text(s, 0.95, 4.78, 7.2, 0.5,
          [("Digital gate passes, verified in seconds.", 15, False, LILAC, BODY)])

names = card(s, 0.95, 5.95, 8.1, 1.05, INK_LIFT, radius=0.1)
t4 = text(s, 1.35, 6.16, 7.4, 0.7,
          [("Sanjay Verma  ·  Palash Shende  ·  Omkar Velonde  ·  Mukul Sharma  ·  Tushar Shinde",
            12.5, True, WHITE, BODY),
           ("Team lead: Sanjay Verma      C-DAC  ·  Final Evaluation, August 2026",
            11, False, MUTED_DARK, BODY)],
          space=6)

timeline[0] = [ids(t1), ids(t2, t3), ids(names, t4)]

s.notes_slide.notes_text_frame.text = (
    "Ten seconds. Name, one line on what it is, move on. The deck is not the "
    "pitch - the demo is."
)


# ========================================================= 2 - the problem
s = blank(prs)
fill_bg(s, WHITE)

h1 = text(s, 0.9, 0.62, 11.5, 0.8,
          [("A gate pass is still a piece of paper", 38, True, INK, HEAD)])
h2 = text(s, 0.9, 1.42, 7.0, 0.4,
          [("Campus entry runs on handwriting and trust.", 15, False, MUTED, BODY)])

problems = [
    ("It cannot be verified",
     "A paper pass can be copied, lent or forged, and nobody at the gate can tell."),
    ("Approval is a phone call",
     "A visitor arrives, a guard rings a department, someone vouches, a name goes in a register."),
    ("Nothing is searchable",
     "Who was on campus last Tuesday at three? The answer is in a book, if it was written down."),
    ("Onboarding is manual",
     "A new cohort means typing hundreds of forms by hand, one account at a time."),
]

groups = [ids(h1, h2)]
y = 2.35
for i, (head, body) in enumerate(problems, 1):
    n = text(s, 0.9, y - 0.06, 0.75, 0.5, [("0%d" % i, 26, True, SOFT, HEAD)])
    a = text(s, 1.62, y, 5.4, 0.32, [(head, 17, True, TEXT, BODY)])
    b = text(s, 1.62, y + 0.36, 5.4, 0.6, [(body, 12.5, False, MUTED, BODY)],
             line_spacing=1.15)
    # Number, heading and description are ONE beat. Fading them separately
    # made a four-item list read as twelve events.
    groups.append(ids(n, a, b))
    y += 1.12

pc = card(s, 7.55, 2.35, 4.85, 4.15, INK, radius=0.05, shadow=True)
mark = rect(s, 8.15, 2.95, 0.55, 0.055, AMBER)
p1 = text(s, 8.15, 3.22, 3.65, 0.35, [("Every one of these", 15, False, LILAC, BODY)])
p2 = text(s, 8.15, 3.66, 3.8, 1.7,
          [("is a checking\nproblem, not a\npaper problem.", 27, True, WHITE, HEAD)],
          line_spacing=1.15)
p3 = text(s, 8.15, 5.55, 3.65, 0.8,
          [("Digitising the form changes nothing.\nThe gate still cannot verify anyone.",
            12, False, MUTED_DARK, BODY)], line_spacing=1.25)

timeline[1] = groups + [ids(pc, mark, p1, p2, p3)]

s.notes_slide.notes_text_frame.text = (
    "About 50 seconds. Do not read the four points - land the last one: the "
    "problem is verification, not paperwork. That is why a PDF pass or a "
    "Google Form would not have solved it."
)


# ======================================================== 3 - the solution
s = blank(prs)
fill_bg(s, WHITE)

h1 = text(s, 0.9, 0.62, 11.5, 0.8,
          [("One system, four people", 38, True, INK, HEAD)])
h2 = text(s, 0.9, 1.42, 9.5, 0.4,
          [("An encrypted QR that a guard can trust in under a second.",
            15, False, MUTED, BODY)])

roles = [
    ("Student & Faculty",
     "Verified profile, photo on file, a standing pass on their phone.",
     VIOLET_MID, glyph_person),
    ("Visitor",
     "Applies online, faculty approves, pass arrives by email before they travel.",
     VIOLET, glyph_door),
    ("Guard",
     "Scans at the gate. One-second verdict, with the holder's photo to check the face.",
     RGBColor(0x7C, 0x3A, 0xED), glyph_scan),
    ("Administration",
     "Every scan recorded and searchable. Entry-only by design - no exit scan to forget.",
     RGBColor(0x59, 0x21, 0xA8), glyph_log),
]

groups = [ids(h1, h2)]
positions = [(0.9, 2.15), (6.75, 2.15), (0.9, 4.05), (6.75, 4.05)]
for (rx, ry), (head, body, tint, glyph) in zip(positions, roles):
    c = card(s, rx, ry, 5.65, 1.72, PAPER, radius=0.06, line=SOFT, shadow=True)
    disc = oval(s, rx + 0.38, ry + 0.48, 0.44, 0.44, tint)
    gl = glyph(s, rx + 0.60, ry + 0.70, WHITE)
    a = text(s, rx + 1.12, ry + 0.40, 4.15, 0.35, [(head, 17, True, INK, BODY)])
    b = text(s, rx + 1.12, ry + 0.82, 4.15, 0.7, [(body, 12.5, False, MUTED, BODY)],
             line_spacing=1.15)
    # The card, its icon and its words are one object as far as the audience
    # is concerned, so they arrive together. An icon that appears before its
    # card is floating on the background.
    groups.append(ids(c, disc, gl, a, b))

bar = card(s, 0.9, 5.95, 11.5, 1.15, INK, radius=0.09, shadow=True)
bmark = rect(s, 1.45, 6.22, 0.055, 0.62, AMBER)
bt = text(s, 1.7, 6.20, 10.2, 0.7,
          [("A whole cohort, from one upload.", 17, True, AMBER_SOFT, BODY),
           ("Students fill in a form. Faculty upload the responses: accounts created, "
            "details verified, photos pulled, passes issued and emails sent.",
            12.5, False, MUTED_DARK, BODY)],
          space=5, line_spacing=1.15)

timeline[2] = groups + [ids(bar, bmark, bt)]

s.notes_slide.notes_text_frame.text = (
    "About 70 seconds. The four cards are a glance. Spend the time on the bar "
    "at the bottom - bulk onboarding is the thing no one else will have built, "
    "and it is the first thing to show in the demo."
)


# ====================================================== 4 - how it is built
s = blank(prs)
fill_bg(s, WHITE)

h1 = text(s, 0.9, 0.62, 11.5, 0.8,
          [("Six services, one gate", 38, True, INK, HEAD)])
h2 = text(s, 0.9, 1.42, 10.5, 0.4,
          [("Spring Boot microservices, a database each, talking over REST and a message queue.",
            15, False, MUTED, BODY)])

services = [
    ("Auth", "Accounts,\nJWT, OTP"),
    ("User", "Profiles,\nverification"),
    ("Gate Pass", "Passes,\nvisitors"),
    ("QR", "AES-256\ntokens, PDFs"),
    ("Campus", "Gates,\npolicy"),
    ("Guard", "Scanning,\nentry log"),
]

groups = [ids(h1, h2)]
row = []
x = 0.9
for name, detail in services:
    c = card(s, x, 2.15, 1.79, 1.8, PAPER, radius=0.08, line=SOFT, shadow=True)
    d = oval(s, x + 0.2, 2.42, 0.2, 0.2, VIOLET_MID)
    a = text(s, x + 0.2, 2.76, 1.5, 0.3, [(name, 15, True, VIOLET, BODY)])
    b = text(s, x + 0.2, 3.14, 1.45, 0.75, [(detail, 11, False, MUTED, BODY)],
             line_spacing=1.15)
    row += ids(c, d, a, b)
    x += 1.93
# All six at once. They are one idea - "six services" - not six ideas, and
# revealing them one by one invites the audience to read rather than listen.
groups.append(row)

infra = card(s, 0.9, 4.22, 11.5, 0.7, INK, radius=0.14, shadow=True)
it = text(s, 1.35, 4.44, 10.6, 0.3,
          [("PostgreSQL  ·  MongoDB  ·  RabbitMQ  ·  Redis  ·  Eureka  ·  Docker  ·  React + TypeScript",
            12.5, True, WHITE, BODY)])
groups.append(ids(infra, it))

decisions = [
    ("Entry only",
     "The system records arrivals. No exit scan means no half-open record when someone forgets."),
    ("A pass is a token",
     "AES-256 encrypted, tied to one holder. A screenshot is useless once the pass is revoked."),
    ("Database per service",
     "No service reads another's tables. They ask over REST, so one can fail without taking the gate down."),
]

trio = []
x = 0.9
for head, body in decisions:
    m = rect(s, x, 5.30, 0.22, 0.055, AMBER)
    a = text(s, x, 5.52, 3.55, 0.3, [(head, 15, True, INK, BODY)])
    b = text(s, x, 5.92, 3.55, 1.1, [(body, 12, False, MUTED, BODY)], line_spacing=1.2)
    trio += ids(m, a, b)
    x += 3.97
groups.append(trio)

timeline[3] = groups

s.notes_slide.notes_text_frame.text = (
    "Do not present this slide. It exists for Q&A - when someone asks why six "
    "services, or what happens when one is down, put it up and answer from it."
)


# ==================================================== 5 - status and future
s = blank(prs)
fill_bg(s, INK)
oval(s, 11.55, -2.65, 4.5, 4.5, INK_LIFT)

h1 = text(s, 0.9, 0.62, 11.0, 0.8,
          [("Working today, and what comes next", 36, True, WHITE, HEAD)])

working = [
    "Bulk onboarding end to end: form to account to pass to email",
    "Encrypted QR passes, issued and delivered as PDF",
    "Faculty approval for visitors, with email at every step",
    "Guard scanning with the holder's photo on screen",
    "Entry log, searchable and paged, for every scan",
]
nexts = [
    "Offline scanning, for when the gate loses network",
    "Face match against the stored photo at the moment of entry",
    "Multi-campus: already in the data model, not yet in the UI",
    "A native guard app instead of a browser",
    "Single sign-on with the college identity provider",
]

groups = [ids(h1)]

# One column per beat. Five bullets arriving individually is ten seconds of
# the audience waiting for a list they could have read in two.
lt = text(s, 0.9, 1.72, 5.4, 0.35, [("WORKING TODAY", 13, True, AMBER, BODY)])
left = ids(lt)
y = 2.28
for item in working:
    d = oval(s, 0.93, y + 0.05, 0.15, 0.15, AMBER)
    t = text(s, 1.30, y, 4.9, 0.5, [(item, 13, False, WHITE, BODY)], line_spacing=1.15)
    left += ids(d, t)
    y += 0.78
groups.append(left)

rt = text(s, 7.0, 1.72, 5.4, 0.35, [("WHAT COMES NEXT", 13, True, LILAC, BODY)])
right = ids(rt)
y = 2.28
for item in nexts:
    d = oval(s, 7.03, y + 0.05, 0.15, 0.15, VIOLET_MID)
    t = text(s, 7.40, y, 4.9, 0.5, [(item, 13, False, SOFT, BODY)], line_spacing=1.15)
    right += ids(d, t)
    y += 0.78
groups.append(right)

cl = text(s, 0.9, 6.62, 11.5, 0.5,
          [("Built to work at any campus - no institution, department or email domain "
            "is written into the code.", 13, True, MUTED_DARK, BODY)])
groups.append(ids(cl))
timeline[4] = groups

s.notes_slide.notes_text_frame.text = (
    "Closing slide, and the one to leave on screen during Q&A. The future-scope "
    "column is what they will push on - offline scanning and face match are the "
    "two with real engineering behind them, so lead with those."
)


prs.save(OUT)


# ------------------------------------------------------- animation & transitions
# python-pptx has no API for either, so the timing tree is written straight
# into each slide part.
#
# ==========================================================================
# GROUPS, NOT SHAPES
# ==========================================================================
# timeline[slide] is a list of GROUPS, and every shape in a group fades in at
# the same instant. The first version animated each shape on its own, so a
# heading and its description arrived as two separate events and a four-item
# list played as twelve - which is what makes an animated deck feel amateur.
#
# The whole slide runs from ONE click: the first group is the click, each
# later group follows the previous automatically. A presenter should never
# have to remember how many times to press the button mid-sentence.
#
#   first shape of the first group   -> clickEffect   (waits for the click)
#   first shape of a later group     -> afterEffect   (follows the previous)
#   every other shape in any group   -> withEffect    (simultaneous)

EFFECT = """<p:par><p:cTn id="{a}" fill="hold"><p:stCondLst><p:cond delay="{delay}"/></p:stCondLst>\
<p:childTnLst><p:par><p:cTn id="{b}" fill="hold"><p:stCondLst><p:cond delay="0"/></p:stCondLst>\
<p:childTnLst><p:par><p:cTn id="{c}" presetID="10" presetClass="entr" presetSubtype="0" \
fill="hold" grpId="0" nodeType="{node}"><p:stCondLst><p:cond delay="0"/></p:stCondLst>\
<p:childTnLst><p:set><p:cBhvr><p:cTn id="{d}" dur="1" fill="hold"><p:stCondLst>\
<p:cond delay="0"/></p:stCondLst></p:cTn><p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl>\
<p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst></p:cBhvr>\
<p:to><p:strVal val="visible"/></p:to></p:set><p:animEffect transition="in" filter="fade">\
<p:cBhvr><p:cTn id="{e}" dur="{dur}"/><p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl></p:cBhvr>\
</p:animEffect></p:childTnLst></p:cTn></p:par></p:childTnLst></p:cTn></p:par></p:childTnLst>\
</p:cTn></p:par>"""

TIMING = """<p:timing><p:tnLst><p:par><p:cTn id="1" dur="indefinite" restart="never" \
nodeType="tmRoot"><p:childTnLst><p:seq concurrent="1" nextAc="seek"><p:cTn id="2" \
dur="indefinite" nodeType="mainSeq"><p:childTnLst>{effects}</p:childTnLst></p:cTn>\
<p:prevCondLst><p:cond evt="onPrev" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond>\
</p:prevCondLst><p:nextCondLst><p:cond evt="onNext" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl>\
</p:cond></p:nextCondLst></p:seq></p:childTnLst></p:cTn></p:par></p:tnLst></p:timing>"""

TRANSITION = '<p:transition spd="med"><p:fade/></p:transition>'

FADE_MS = 450        # long enough to read as a fade, short enough not to wait
GROUP_GAP_MS = 200   # between groups


def slide_timing(groups):
    effects, uid = [], 3
    first_group = True

    for group in groups:
        for i, spid in enumerate(group):
            if i:
                node, delay = "withEffect", 0
            elif first_group:
                node, delay = "clickEffect", 0
            else:
                node, delay = "afterEffect", GROUP_GAP_MS

            effects.append(EFFECT.format(
                a=uid, b=uid + 1, c=uid + 2, d=uid + 3, e=uid + 4,
                spid=spid, node=node, delay=delay, dur=FADE_MS,
            ))
            uid += 5
        first_group = False

    return TIMING.format(effects="".join(effects)) if effects else ""


def animate(path):
    src = Path(path)
    tmp = src.with_suffix(".tmp.pptx")

    with zipfile.ZipFile(src) as zin, zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)

            m = re.fullmatch(r"ppt/slides/slide(\d+)\.xml", item.filename)
            if m:
                groups = timeline.get(int(m.group(1)) - 1, [])
                xml = data.decode("utf-8")
                xml = xml.replace("</p:sld>", TRANSITION + slide_timing(groups) + "</p:sld>")
                data = xml.encode("utf-8")

            zout.writestr(item, data)

    shutil.move(str(tmp), str(src))


animate(OUT)
print("%s  -  %d slides, %d animation steps"
      % (OUT, len(prs.slides._sldIdLst), sum(len(g) for g in timeline.values())))
